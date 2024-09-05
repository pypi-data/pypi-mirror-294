import PyPDF2
from pptx import Presentation
from docx import Document
import re
import os
import argparse
import logging
from PyPDF2 import PdfReader 

logging.basicConfig(level=logging.WARNING)

def extract_pdf_to_markdown(file_path):
    # Make sure PdfReader is used in this function
    reader = PdfReader(file_path)
    markdown = f"# {os.path.basename(file_path)}\n\n"
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        markdown += f"## Page {i+1}\n\n{text}\n\n"
    return markdown

def extract_pptx_to_markdown(file_path):
    prs = Presentation(file_path)
    markdown = f"# {os.path.basename(file_path)}\n\n"
    for i, slide in enumerate(prs.slides):
        markdown += f"## Slide {i+1}\n\n"
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                markdown += f"{shape.text}\n\n"
    return markdown

def extract_docx_to_markdown(file_path):
    doc = Document(file_path)
    markdown = f"# {os.path.basename(file_path)}\n\n"
    for para in doc.paragraphs:
        if para.style.name.startswith('Heading'):
            level = para.style.name.split()[-1]
            level = int(level) if level.isdigit() else 1
            markdown += f"{'#' * level} {para.text}\n\n"
        else:
            markdown += f"{para.text}\n\n"
    return markdown

def clean_markdown(markdown_text):
    # Normalize newlines
    cleaned = re.sub(r'\r\n|\r', '\n', markdown_text)
    
    # Replace multiple newlines with double newlines
    cleaned = re.sub(r'\n{3,}', '\n\n', cleaned)
    
    # Remove extra spaces within lines
    cleaned = re.sub(r'[ \t]+', ' ', cleaned)
    
    # Remove leading/trailing whitespace from each line
    cleaned = '\n'.join(line.strip() for line in cleaned.split('\n'))
    
    # Remove leading/trailing whitespace from the entire text
    cleaned = cleaned.strip()
    
    return cleaned

def process_file_to_markdown(file_path, output_dir):
    _, ext = os.path.splitext(file_path)
    if ext.lower() == '.pdf':
        markdown = extract_pdf_to_markdown(file_path)
    elif ext.lower() == '.pptx':
        markdown = extract_pptx_to_markdown(file_path)
    elif ext.lower() == '.docx':
        markdown = extract_docx_to_markdown(file_path)
    else:
        print(f"Skipping unsupported file: {file_path}")
        return

    markdown = clean_markdown(markdown)
    
    rel_path = os.path.relpath(file_path, start=os.path.dirname(file_path))
    output_path = os.path.join(output_dir, os.path.splitext(rel_path)[0] + '.md')
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(markdown)
    
    logging.warning(f"Markdown file created: {output_path}")

def process_folder(input_dir, output_dir):
    for root, _, files in os.walk(input_dir):
        for file in files:
            file_path = os.path.join(root, file)
            process_file_to_markdown(file_path, output_dir)

def main():
    parser = argparse.ArgumentParser(
        description="Convert PDF, PPTX, and DOCX files to Markdown format.",
        epilog="This script processes all compatible files in the input directory and its subdirectories, "
               "creating a corresponding structure of Markdown files in the output directory."
    )
    parser.add_argument('input_dir', help="Path to the input directory containing the documents to be converted")
    parser.add_argument('output_dir', help="Path to the output directory where Markdown files will be saved")
    parser.add_argument('-v', '--verbose', action='store_true', help="Increase output verbosity")

    global args
    args = parser.parse_args()

    if args.verbose:
        print(f"Processing files from: {args.input_dir}")
        print(f"Saving Markdown files to: {args.output_dir}")

    process_folder(args.input_dir, args.output_dir)

    if args.verbose:
        print("Processing complete.")

if __name__ == "__main__":
    main()
